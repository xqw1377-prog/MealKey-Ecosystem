from __future__ import annotations

import base64
import csv
import io
import json
import mimetypes
import os
import re
from dataclasses import asdict, dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from zipfile import BadZipFile

from fastapi import UploadFile

from app.services.llm_engine.bindings import PurposeModelCandidate, resolve_candidate_api_key
from app.services.llm_engine.client import ChatMessage, chat_completion
from app.services.speech_service import transcribe_audio_bytes

MAX_ATTACHMENTS = 12
MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_TOTAL_BYTES = 50 * 1024 * 1024
MAX_EXTRACTED_CHARS = 12000
MAX_CONTEXT_CHARS = 18000

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".jsonl",
    ".csv",
    ".tsv",
    ".log",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
    ".conf",
    ".env",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".css",
    ".scss",
    ".less",
    ".sql",
    ".sh",
    ".ps1",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac", ".webm"}


@dataclass
class ParsedAttachment:
    name: str
    content_type: str
    size_bytes: int
    extension: str
    kind: str
    summary: str
    extracted_text: str = ""
    warnings: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_public_dict(self) -> dict[str, Any]:
        data = asdict(self)
        data["size_kb"] = round(self.size_bytes / 1024, 1)
        return data


class _HtmlTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = (data or "").strip()
        if text:
            self.parts.append(text)

    def get_text(self) -> str:
        return " ".join(self.parts)


def _truncate(text: str, limit: int = MAX_EXTRACTED_CHARS) -> str:
    clean = re.sub(r"\n{3,}", "\n\n", (text or "").strip())
    if len(clean) <= limit:
        return clean
    return clean[:limit].rstrip() + "\n…"


def _preview_summary(text: str, *, limit: int = 180) -> str:
    if not text:
        return ""
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines() if line.strip()]
    preview = "；".join(lines[:3]) if lines else re.sub(r"\s+", " ", text).strip()
    if len(preview) <= limit:
        return preview
    return preview[:limit].rstrip("，；。 ") + "…"


def _safe_filename(name: str | None) -> str:
    return Path(name or "attachment").name or "attachment"


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _text_from_html(raw: str) -> str:
    parser = _HtmlTextExtractor()
    parser.feed(raw)
    return parser.get_text()


def _sheet_to_text(rows: list[list[Any]], *, max_rows: int = 20, max_cols: int = 10) -> str:
    out: list[str] = []
    for row in rows[:max_rows]:
        cells = [str(cell).strip() for cell in row[:max_cols] if cell not in (None, "")]
        if cells:
            out.append(" | ".join(cells))
    return "\n".join(out)


def _parse_csv_text(raw: str, delimiter: str = ",") -> str:
    reader = csv.reader(io.StringIO(raw), delimiter=delimiter)
    rows = [row for _, row in zip(range(20), reader)]
    return _sheet_to_text(rows)


def _resolve_vision_candidate() -> PurposeModelCandidate | None:
    api_key = (os.getenv("QWEN_API_KEY") or os.getenv("DASHSCOPE_API_KEY") or "").strip()
    if not api_key:
        return None
    return PurposeModelCandidate(
        id="vision:qwen-vl",
        provider="qianwen",
        model=(os.getenv("QWEN_VISION_MODEL") or "qwen-vl-max-latest").strip(),
        base_url=(os.getenv("QWEN_BASE_URL") or "https://dashscope.aliyuncs.com/compatible-mode/v1").strip(),
        api_key_env="QWEN_API_KEY" if os.getenv("QWEN_API_KEY") else "DASHSCOPE_API_KEY",
        quality_tier="balanced",
    )


def _transcribe_audio(data: bytes, filename: str, content_type: str) -> tuple[str, list[str]]:
    return transcribe_audio_bytes(data, filename, content_type)


def _analyze_image(data: bytes, filename: str, content_type: str) -> tuple[str, list[str]]:
    candidate = _resolve_vision_candidate()
    if not candidate:
        return "", ["当前未配置图片理解模型，已保留图片文件元信息。"]
    api_key = resolve_candidate_api_key(candidate)
    if not api_key:
        return "", ["当前未找到图片理解模型的可用密钥。"]
    data_uri = f"data:{content_type or 'image/png'};base64,{base64.b64encode(data).decode('ascii')}"
    prompt = (
        "请先识别图片中的主要内容、可见文字和关键信息。"
        "如果这是经营截图、表格、菜单、发票、聊天记录或商品图，请优先提取与经营决策有关的要点。"
        "输出简洁中文摘要，控制在 180 字内。"
    )
    result = chat_completion(
        api_key=api_key,
        base_url=candidate.base_url,
        model=candidate.model,
        messages=[
            ChatMessage(
                role="user",
                content=[
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_uri}},
                ],
            )
        ],
        temperature=0.2,
        max_tokens=400,
        timeout_seconds=90,
    )
    return _truncate(result.content, 1000), []


def _parse_pdf(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    warnings: list[str] = []
    try:
        from pypdf import PdfReader
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"PDF 解析库不可用：{exc}"]
    try:
        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for page in reader.pages[:8]:
            pages.append(page.extract_text() or "")
        return _truncate("\n\n".join(pages)), {"pages": len(reader.pages)}, warnings
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"PDF 解析失败：{exc}"]


def _parse_docx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    try:
        from docx import Document
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"DOCX 解析库不可用：{exc}"]
    try:
        doc = Document(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return _truncate("\n".join(paras)), {"paragraphs": len(paras)}, []
    except BadZipFile as exc:
        return "", {}, [f"DOCX 文件损坏：{exc}"]
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"DOCX 解析失败：{exc}"]


def _parse_pptx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"PPTX 解析库不可用：{exc}"]
    try:
        presentation = Presentation(io.BytesIO(data))
        texts: list[str] = []
        for slide in presentation.slides[:15]:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    texts.append(text.strip())
        return _truncate("\n".join(texts)), {"slides": len(presentation.slides)}, []
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"PPTX 解析失败：{exc}"]


def _parse_xlsx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"表格解析库不可用：{exc}"]
    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        text_blocks: list[str] = []
        sheet_summaries: list[dict[str, Any]] = []
        for sheet in workbook.worksheets[:5]:
            rows = []
            for row in sheet.iter_rows(min_row=1, max_row=20, max_col=10, values_only=True):
                rows.append(list(row))
            sheet_text = _sheet_to_text(rows)
            if sheet_text:
                text_blocks.append(f"[工作表] {sheet.title}\n{sheet_text}")
            sheet_summaries.append({"sheet": sheet.title, "rows_previewed": len(rows)})
        return _truncate("\n\n".join(text_blocks)), {"sheets": sheet_summaries}, []
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"表格解析失败：{exc}"]


def _parse_xls(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    try:
        import xlrd
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"XLS 解析库不可用：{exc}"]
    try:
        book = xlrd.open_workbook(file_contents=data)
        text_blocks: list[str] = []
        sheet_summaries: list[dict[str, Any]] = []
        for sheet in book.sheets()[:5]:
            rows = [sheet.row_values(i, 0, min(sheet.ncols, 10)) for i in range(min(sheet.nrows, 20))]
            sheet_text = _sheet_to_text(rows)
            if sheet_text:
                text_blocks.append(f"[工作表] {sheet.name}\n{sheet_text}")
            sheet_summaries.append({"sheet": sheet.name, "rows_previewed": min(sheet.nrows, 20)})
        return _truncate("\n\n".join(text_blocks)), {"sheets": sheet_summaries}, []
    except Exception as exc:  # noqa: BLE001
        return "", {}, [f"XLS 解析失败：{exc}"]


def _parse_plain_text(data: bytes, extension: str, content_type: str) -> tuple[str, dict[str, Any], list[str]]:
    raw = _decode_text(data)
    text = raw
    metadata: dict[str, Any] = {}
    if extension == ".json":
        try:
            parsed = json.loads(raw)
            metadata["json_type"] = type(parsed).__name__
            text = json.dumps(parsed, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            pass
    elif extension in {".csv", ".tsv"}:
        text = _parse_csv_text(raw, delimiter="\t" if extension == ".tsv" else ",")
    elif extension in {".html", ".htm"} or "html" in content_type:
        text = _text_from_html(raw)
    return _truncate(text), metadata, []


def parse_attachment_bytes(
    *,
    filename: str,
    content_type: str,
    data: bytes,
) -> ParsedAttachment:
    name = _safe_filename(filename)
    extension = Path(name).suffix.lower()
    inferred_type = content_type or mimetypes.guess_type(name)[0] or "application/octet-stream"
    size_bytes = len(data)
    warnings: list[str] = []
    metadata: dict[str, Any] = {}
    extracted_text = ""
    kind = "binary"

    if size_bytes > MAX_FILE_BYTES:
        return ParsedAttachment(
            name=name,
            content_type=inferred_type,
            size_bytes=size_bytes,
            extension=extension,
            kind="oversize",
            summary=f"{name} 超过单文件大小限制，当前仅保留元信息。",
            warnings=[f"单文件上限 {MAX_FILE_BYTES // (1024 * 1024)}MB。"],
        )

    if extension in TEXT_EXTENSIONS or inferred_type.startswith("text/") or inferred_type in {"application/json", "application/xml"}:
        kind = "text"
        extracted_text, metadata, warnings = _parse_plain_text(data, extension, inferred_type)
    elif extension == ".pdf" or inferred_type == "application/pdf":
        kind = "pdf"
        extracted_text, metadata, warnings = _parse_pdf(data)
    elif extension == ".docx" or inferred_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        kind = "document"
        extracted_text, metadata, warnings = _parse_docx(data)
    elif extension in {".xlsx", ".xlsm"} or inferred_type.startswith("application/vnd.openxmlformats-officedocument.spreadsheetml"):
        kind = "spreadsheet"
        extracted_text, metadata, warnings = _parse_xlsx(data)
    elif extension == ".xls":
        kind = "spreadsheet"
        extracted_text, metadata, warnings = _parse_xls(data)
    elif extension == ".pptx" or inferred_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        kind = "presentation"
        extracted_text, metadata, warnings = _parse_pptx(data)
    elif extension in IMAGE_EXTENSIONS or inferred_type.startswith("image/"):
        kind = "image"
        extracted_text, warnings = _analyze_image(data, name, inferred_type)
        try:
            from PIL import Image

            image = Image.open(io.BytesIO(data))
            metadata["width"] = image.width
            metadata["height"] = image.height
        except Exception:  # noqa: BLE001
            pass
    elif extension in AUDIO_EXTENSIONS or inferred_type.startswith("audio/"):
        kind = "audio"
        extracted_text, warnings = _transcribe_audio(data, name, inferred_type)
    else:
        try:
            extracted_text = _truncate(_decode_text(data))
            kind = "text"
            warnings.append("文件类型未显式识别，已按文本尝试解析。")
        except Exception:  # noqa: BLE001
            warnings.append("暂不支持直接提取该文件类型内容，已保留文件元信息。")

    line_count = extracted_text.count("\n") + 1 if extracted_text else 0
    metadata["line_count"] = line_count
    summary_bits = [f"{name}（{kind}）"]
    if metadata.get("pages"):
        summary_bits.append(f"{metadata['pages']} 页")
    if metadata.get("slides"):
        summary_bits.append(f"{metadata['slides']} 页幻灯片")
    if metadata.get("width") and metadata.get("height"):
        summary_bits.append(f"{metadata['width']}×{metadata['height']}")
    if extracted_text:
        summary_bits.append(_preview_summary(extracted_text))
        if len(extracted_text) >= MAX_EXTRACTED_CHARS - 4:
            summary_bits.append("内容较长，已截取重点片段")
    elif warnings:
        summary_bits.append(warnings[0])
    else:
        summary_bits.append("已接收文件。")

    return ParsedAttachment(
        name=name,
        content_type=inferred_type,
        size_bytes=size_bytes,
        extension=extension,
        kind=kind,
        summary=" · ".join(summary_bits),
        extracted_text=extracted_text,
        warnings=warnings,
        metadata=metadata,
    )


async def parse_upload_files(files: list[UploadFile]) -> list[ParsedAttachment]:
    if len(files) > MAX_ATTACHMENTS:
        raise ValueError(f"单次最多上传 {MAX_ATTACHMENTS} 个文件。")
    total_bytes = 0
    parsed: list[ParsedAttachment] = []
    for upload in files:
        data = await upload.read()
        total_bytes += len(data)
        if total_bytes > MAX_TOTAL_BYTES:
            raise ValueError(f"单次上传总大小不能超过 {MAX_TOTAL_BYTES // (1024 * 1024)}MB。")
        parsed.append(
            parse_attachment_bytes(
                filename=upload.filename or "attachment",
                content_type=upload.content_type or "",
                data=data,
            )
        )
        await upload.close()
    return parsed


def build_attachment_context(parsed_files: list[ParsedAttachment]) -> str:
    if not parsed_files:
        return ""
    sections = ["用户补充了以下附件，请优先结合附件内容作答："]
    for index, item in enumerate(parsed_files, start=1):
        section = [f"{index}. {item.summary}"]
        if item.warnings:
            section.append(f"注意：{'；'.join(item.warnings[:2])}")
        if item.extracted_text:
            section.append(f"提取内容：\n{_truncate(item.extracted_text, 4000)}")
        sections.append("\n".join(section))
    return _truncate("\n\n".join(sections), MAX_CONTEXT_CHARS)


# ═══════════════════════════════════════════════════════════
# 文件路径读取——给路径→读取→蒸馏→接入对话
# ═══════════════════════════════════════════════════════════


# P0-1 安全：允许读取的目录白名单
_ALLOWED_READ_DIRS = [
    Path("uploads").resolve(),
    Path("data").resolve(),
    Path("static").resolve(),
]

# P0-1 安全：禁止读取的扩展名
_FORBIDDEN_EXTENSIONS = {".env", ".db", ".sqlite", ".sqlite3", ".pem", ".key", ".pfx", ".p12"}

# P0-1 安全：允许的扩展名白名单
_ALLOWED_READ_EXTENSIONS = {
    ".txt", ".md", ".csv", ".xlsx", ".xls", ".pdf", ".json", ".xml", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".gif", ".webp",
}


def _is_path_safe(path: Path) -> tuple[bool, str]:
    """检查路径是否安全可读（沙箱校验）。"""
    try:
        resolved = path.resolve()
    except Exception:
        return False, "路径解析失败"

    # 禁止敏感扩展名
    if resolved.suffix.lower() in _FORBIDDEN_EXTENSIONS:
        return False, f"安全限制：不允许读取 {resolved.suffix} 文件"

    # 只允许白名单扩展名
    if resolved.suffix.lower() not in _ALLOWED_READ_EXTENSIONS:
        return False, f"不支持的文件类型: {resolved.suffix}"

    # 检查是否在允许的目录内
    for allowed_dir in _ALLOWED_READ_DIRS:
        try:
            resolved.relative_to(allowed_dir)
            return True, ""
        except ValueError:
            continue

    # 如果允许的目录不存在，也拒绝
    return False, "安全限制：文件必须在 uploads/ data/ static/ 目录内"


def read_file_by_path(file_path: str) -> dict[str, Any]:
    """从服务端文件路径直接读取并解析（蒸馏）。

    老板可以给一个文件路径（如成本表/菜单图/报表），
    系统直接读取→解析→蒸馏成文本→接入对话上下文。

    支持：.txt .md .csv .xlsx .pdf .json .xml .html
    安全：限制在 uploads/ data/ static/ 目录内，禁止 .env/.db 等敏感文件。
    """
    path = Path(file_path.strip().strip('"').strip("'"))

    # P0-1 安全沙箱校验
    is_safe, safety_reason = _is_path_safe(path)
    if not is_safe:
        return {"ok": False, "error": safety_reason}

    if not path.exists():
        return {"ok": False, "error": f"文件不存在: {file_path}"}
    if not path.is_file():
        return {"ok": False, "error": f"不是文件: {file_path}"}
    if path.stat().st_size > MAX_FILE_BYTES:
        return {"ok": False, "error": f"文件超过 {MAX_FILE_BYTES // 1024 // 1024}MB 限制"}

    ext = path.suffix.lower()
    data = path.read_bytes()
    warnings: list[str] = []

    extracted_text = ""
    structured: dict[str, Any] = {}

    try:
        if ext in TEXT_EXTENSIONS:
            raw = data.decode("utf-8", errors="replace")
            extracted_text = _truncate(raw, MAX_EXTRACTED_CHARS)
        elif ext == ".csv":
            raw = data.decode("utf-8", errors="replace")
            extracted_text = _parse_csv_text(raw)
        elif ext in (".xlsx", ".xls"):
            extracted_text, structured, warnings = _parse_xlsx(data)
        elif ext == ".pdf":
            extracted_text, structured, warnings = _parse_pdf(data)
        elif ext == ".json":
            raw = data.decode("utf-8", errors="replace")
            parsed = json.loads(raw)
            extracted_text = _truncate(json.dumps(parsed, ensure_ascii=False, indent=2), MAX_EXTRACTED_CHARS)
        elif ext in (".html", ".htm"):
            raw = data.decode("utf-8", errors="replace")
            extracted_text = _truncate(_HTMLTextExtractor().extract_text(raw), MAX_EXTRACTED_CHARS)
        else:
            # 尝试当文本读
            try:
                raw = data.decode("utf-8", errors="replace")
                extracted_text = _truncate(raw[:5000], MAX_EXTRACTED_CHARS)
                warnings.append(f"未知扩展名 {ext}，按文本读取")
            except Exception:
                return {"ok": False, "error": f"不支持文件类型: {ext}"}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"解析失败: {type(exc).__name__}: {str(exc)[:100]}"}

    summary = f"{path.name}（{path.stat().st_size // 1024}KB）"

    return {
        "ok": True,
        "file_path": str(path),
        "file_name": path.name,
        "summary": summary,
        "extracted_text": extracted_text[:8000],
        "structured": structured if structured else None,
        "warnings": warnings,
        "context_for_chat": f"文件 {path.name} 的内容：\n{_truncate(extracted_text, 6000)}",
    }


class _HTMLTextExtractor(HTMLParser):
    """简单的 HTML 文本提取器。"""

    def __init__(self):
        super().__init__()
        self._text: list[str] = []

    def handle_data(self, data: str) -> None:
        self._text.append(data)

    def extract_text(self, html: str) -> str:
        self._text = []
        self.feed(html)
        return " ".join(t.strip() for t in self._text if t.strip())
